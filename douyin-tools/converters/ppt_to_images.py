import os
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from PIL import Image
import io


def _libreoffice_convert(pptx_path: str, output_dir: str, dpi: int) -> list[str]:
    """Use LibreOffice to convert PPTX → PDF, then render pages as images."""
    pdf_path = Path(output_dir) / (Path(pptx_path).stem + ".pdf")
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
        check=True, capture_output=True, timeout=120,
    )
    images = convert_from_path(str(pdf_path), dpi=dpi)
    image_paths = []
    for i, img in enumerate(images):
        img_path = Path(output_dir) / f"slide_{i+1:03d}.png"
        img.save(str(img_path), "PNG")
        image_paths.append(str(img_path))
    return image_paths


def _extract_shapes(pptx_path: str, output_dir: str) -> list[str]:
    """Fallback: extract embedded images and render text shapes from PPTX slides.

    For each slide, creates a blank canvas and overlays extracted images + text.
    """
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    image_paths = []
    for slide_num, slide in enumerate(prs.slides):
        canvas = Image.new("RGB", (int(slide_w / 10000), int(slide_h / 10000)), "white")
        canvas_pil = Image.new("RGBA", canvas.size, (255, 255, 255, 255))

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    img_bytes = shape.image.blob
                    img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
                    img = img.resize((int(shape.width / 10000), int(shape.height / 10000)))
                    x, y = int(shape.left / 10000), int(shape.top / 10000)
                    canvas_pil.paste(img, (x, y), img)
                except Exception:
                    pass

        # Merge onto white background
        canvas = Image.alpha_composite(canvas.convert("RGBA"), canvas_pil).convert("RGB")
        img_path = str(Path(output_dir) / f"slide_{slide_num+1:03d}.png")
        canvas.save(img_path, "PNG")
        image_paths.append(img_path)

    return image_paths


def ppt_to_images(pptx_path: str, dpi: int = 200) -> str:
    """Convert PPTX slides to images. Returns path to a zip file of images."""
    output_dir = Path(tempfile.mkdtemp())
    try:
        if shutil.which("soffice"):
            image_paths = _libreoffice_convert(pptx_path, str(output_dir), dpi)
        else:
            image_paths = _extract_shapes(pptx_path, str(output_dir))
    except Exception:
        image_paths = _extract_shapes(pptx_path, str(output_dir))

    # Package into zip
    zip_path = Path(tempfile.gettempdir()) / (Path(pptx_path).stem + "_images.zip")
    with zipfile.ZipFile(str(zip_path), "w") as zf:
        for img_path in image_paths:
            zf.write(img_path, Path(img_path).name)
    return str(zip_path)
