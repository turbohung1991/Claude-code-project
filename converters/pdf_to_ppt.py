import os
import tempfile
from pathlib import Path
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def pdf_to_ppt(pdf_path: str, dpi: int = 200) -> str:
    """Convert PDF to PPTX. Each page becomes a full-slide image."""
    images = convert_from_path(pdf_path, dpi=dpi)
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            img.save(tmp.name, "PNG")
            slide.shapes.add_picture(tmp.name, 0, 0, slide_width, slide_height)

    output_path = Path(pdf_path).with_suffix(".pptx")
    prs.save(output_path)

    # Clean up temp files
    for img_file in Path(tempfile.gettempdir()).glob("*.png"):
        try:
            img_file.unlink()
        except OSError:
            pass

    return str(output_path)
