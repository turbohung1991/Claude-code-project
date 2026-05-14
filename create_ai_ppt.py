from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK   = RGBColor(0x0B, 0x14, 0x2D)   # deep navy
BG_CARD   = RGBColor(0x13, 0x1F, 0x42)   # card navy
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # bright blue
ACCENT2   = RGBColor(0x06, 0xB6, 0xD4)   # cyan
ACCENT3   = RGBColor(0x8B, 0x5C, 0xF6)   # purple
GREEN     = RGBColor(0x10, 0xB9, 0x81)   # green
ORANGE    = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED       = RGBColor(0xEF, 0x44, 0x44)   # red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT     = RGBColor(0xE5, 0xE7, 0xEB)
SUBTLE    = RGBColor(0x6B, 0x72, 0x80)

# ── Helper functions ──
def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE, line_spacing=1.5, font_name='Microsoft YaHei'):
    """lines is a list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, fs, clr = item, False, None, None
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            fs = item[2] if len(item) > 2 else None
            clr = item[3] if len(item) > 3 else None
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(fs or font_size)
        p.font.color.rgb = clr or color
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_tag(slide, left, top, text, color=ACCENT):
    shape = add_rect(slide, left, top, Inches(1.6), Inches(0.35), color, 0.1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER

def add_section_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(0.7), title, font_size=32, bold=True, color=WHITE)
    # accent bar
    add_rect(slide, Inches(0.8), Inches(1.15), Inches(1.2), Inches(0.06), ACCENT)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.35), Inches(12), Inches(0.5), subtitle, font_size=14, color=GRAY)

def add_card(slide, left, top, width, height, color=None, border_color=None):
    if color is None:
        color = BASE_CARD
    """Add a rounded card background"""
    shape = add_rect(slide, left, top, width, height, color, 0.08)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape

BASE_CARD = RGBColor(0x16, 0x24, 0x4A)

def add_bottom_bar(slide, text="轻量化 AI 赋能策略 | 卫生护理企业"):
    add_rect(slide, Inches(0), Inches(7.15), Inches(13.333), Inches(0.35), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(7.18), Inches(12), Inches(0.3), text, font_size=9, color=WHITE)

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.2), Inches(7.18), Inches(1), Inches(0.3), str(num), font_size=9, color=WHITE, alignment=PP_ALIGN.RIGHT)

# ═══════════════════════════════════
# SLIDE 1 - COVER
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# Large decorative circle top right
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x13, 0x1F, 0x42)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.8), Inches(3), Inches(3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x1A, 0x2A, 0x55)
circle2.line.fill.background()

# Small accent dots
for x, y, c in [(Inches(11.2), Inches(0.3), ACCENT2), (Inches(12.0), Inches(1.8), ACCENT3), (Inches(10.5), Inches(2.2), ACCENT)]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = c
    dot.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.5), "AI 赋能", font_size=20, color=ACCENT2, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.1), Inches(10), Inches(1.5), "轻量化 · 可落地 · 见效快", font_size=46, bold=True, color=WHITE)
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(8), Inches(0.8), "卫生护理企业 AI 赋能落地策略", font_size=22, color=GRAY)

# Separator line
add_rect(slide, Inches(0.8), Inches(4.4), Inches(2.5), Inches(0.05), ACCENT)

add_textbox(slide, Inches(0.8), Inches(5.0), Inches(8), Inches(0.5),
            "不是 AI 项目，是 AI 习惯\n不是技术变革，是工作方式升级\n不建系统、不招团队、不花大钱\n给每个人配一个好用的「实习生」，让效果自己传染",
            font_size=15, color=SUBTLE)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(5), Inches(0.4), "2025年6月", font_size=13, color=GRAY)
add_bottom_bar(slide)

# ═══════════════════════════════════
# SLIDE 2 - TABLE OF CONTENTS
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "目录", "CONTENTS")

items = [
    ("01", "核心理念", "不做大系统，让人人都会用 AI", ACCENT),
    ("02", "按岗位的 AI 使用清单", "跨境电商 · 外贸 · 生产 · 职能 · 管理层", ACCENT2),
    ("03", "推广方式：传染式扩散", "种子选手 → 效果示范 → 模板沉淀 → 自发传播", GREEN),
    ("04", "资源需求（极简版）", "全年不到 3 万元", ORANGE),
    ("05", "注意事项", "数据安全 · 心态 · 工具选择 · 带头示范", ACCENT3),
    ("06", "一张图总结", "", RED),
]

for i, (num, title, desc, color) in enumerate(items):
    y = Inches(1.8) + Inches(0.82) * i
    add_rect(slide, Inches(1.0), y, Inches(0.7), Inches(0.7), color, 0.15)
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(0.7), Inches(0.5), num, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.2), y + Inches(0.05), Inches(5), Inches(0.4), title, font_size=20, bold=True, color=WHITE)
    add_textbox(slide, Inches(2.2), y + Inches(0.42), Inches(8), Inches(0.3), desc, font_size=13, color=GRAY)

add_bottom_bar(slide)
add_page_number(slide, 2)

# ═══════════════════════════════════
# SLIDE 3 - CORE CONCEPT
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "核心理念", "不是搞一个 AI 项目，而是让 AI 成为每个人的工作习惯")

# Three pillar cards
pillars = [
    ("不建系统", "不采购 AI 中台\n不改造现有 IT 架构\n直接用现成工具", ACCENT, "🔧"),
    ("不招团队", "不招算法工程师\n不建数据科学团队\n现有人员技能升级", ACCENT2, "👥"),
    ("不花大钱", "全年预算 < 3 万\nChatGPT Team 账号\n几场分享会 + 茶歇", GREEN, "💰"),
]
for i, (title, desc, color, icon) in enumerate(pillars):
    left = Inches(0.8) + Inches(4.1) * i
    card = add_rect(slide, left, Inches(2.2), Inches(3.7), Inches(3.0), BASE_CARD, 0.1)
    # top accent strip
    add_rect(slide, left, Inches(2.2), Inches(3.7), Inches(0.06), color)
    add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.1), Inches(0.5), title, font_size=24, bold=True, color=color)
    add_textbox(slide, left + Inches(0.3), Inches(3.2), Inches(3.1), Inches(1.5), desc, font_size=15, color=LIGHT)

# Bottom key message
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(12), Inches(0.8),
            "核心理念：把 AI 当成「给每个岗位配一个聪明的实习生」—— 教大家用、给好模板、养成习惯",
            font_size=18, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 3)

# ═══════════════════════════════════
# SLIDE 4 - CROSS-BORDER E-COMMERCE (highest impact)
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "按岗位的 AI 使用清单（一）", "跨境电商运营 — 影响最大，最先推行")
add_tag(slide, Inches(0.8), Inches(1.7), "最高优先级", RED)

scenarios_cb = [
    ("Listing 文案生成", "输入产品参数 + 卖点 → AI 输出英文五点描述 + A+ 文案", "每天省 2-3 小时", ACCENT),
    ("多语言翻译优化", "一键产出日/西/德语 listing，质量够用", "省翻译公司费用", ACCENT2),
    ("差评回复模板", "丢入差评 → AI 生成得体外文回复", "客服不用想了", GREEN),
    ("广告文案 A/B 测试", "一条卖点产出 10 套 FB/Google 广告标题+文案", "投手效率翻倍", ORANGE),
    ("竞品分析", "竞品 listing 导入 → AI 总结卖点策略、定价、用户关注点", "分钟级出报告", ACCENT3),
]

for i, (title, desc, result, color) in enumerate(scenarios_cb):
    y = Inches(2.2) + Inches(0.95) * i
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.12), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(0.45), Inches(0.4), str(i+1), font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.8), y + Inches(0.05), Inches(4), Inches(0.35), title, font_size=18, bold=True, color=WHITE)
    add_textbox(slide, Inches(1.8), y + Inches(0.4), Inches(6), Inches(0.35), desc, font_size=13, color=GRAY)
    add_textbox(slide, Inches(9.0), y + Inches(0.1), Inches(3.5), Inches(0.5), f"效果：{result}", font_size=14, bold=True, color=color)

add_textbox(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
            "落地动作：1 小时实操培训 + 配一套 Prompt 模板库（listing / 广告 / 客服 / 分析各 3-5 个模板）",
            font_size=13, color=ACCENT2)

add_bottom_bar(slide)
add_page_number(slide, 4)

# ═══════════════════════════════════
# SLIDE 5 - FOREIGN TRADE + PRODUCTION
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "按岗位的 AI 使用清单（二）", "外贸/大客户 + 生产供应链")

# Two columns
# Left column - Foreign Trade
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4), "外贸 / 大客户业务", font_size=20, bold=True, color=ACCENT2)
add_rect(slide, Inches(0.8), Inches(2.25), Inches(1.0), Inches(0.04), ACCENT2)

ft_items = [
    "开发信：输入客户背景 → AI 生成针对性英文开发信",
    "合同要点提取：英文合同 PDF → 提取关键条款",
    "邮件辅助：给要点 → 正式商务邮件；长邮件 → 三句话总结",
    "产品报价单：说出需求 → AI 生成格式化英文报价单草稿",
]
for i, item in enumerate(ft_items):
    add_textbox(slide, Inches(0.8), Inches(2.55) + Inches(0.55)*i, Inches(5.5), Inches(0.5),
                f"✦  {item}", font_size=13, color=LIGHT)

# Right column - Production
add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.4), "生产 / 供应链", font_size=20, bold=True, color=GREEN)
add_rect(slide, Inches(7.2), Inches(2.25), Inches(1.0), Inches(0.04), GREEN)

prod_items = [
    "排产方案参考：Excel 计划表 + 约束条件 → AI 给排产建议",
    "异常处理预案：机器故障 → AI 给出调整方案",
    "物料需求估算：排产计划 → AI 算出主要原材料日需求",
    "SOP 润色/翻译：中文 SOP → 工位可用简洁操作卡",
]
for i, item in enumerate(prod_items):
    add_textbox(slide, Inches(7.2), Inches(2.55) + Inches(0.55)*i, Inches(5.5), Inches(0.5),
                f"✦  {item}", font_size=13, color=LIGHT)

# Bottom - Admin & Management
add_rect(slide, Inches(0.8), Inches(5.0), Inches(12), Inches(0.02), RGBColor(0x2A, 0x3A, 0x6A))

add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(0.4), "财务 / 人事 / 行政", font_size=18, bold=True, color=ACCENT3)
admin_items = ["制度问答：员工手册 PDF → AI 秒答差旅标准", "会议纪要：录音转文字 → AI 整理「结论+待办+责任人」", "通知起草：给要点 → AI 写出正式通知"]
for i, item in enumerate(admin_items):
    add_textbox(slide, Inches(0.8), Inches(5.65) + Inches(0.42)*i, Inches(5.5), Inches(0.4), f"▸  {item}", font_size=12, color=LIGHT)

add_textbox(slide, Inches(7.2), Inches(5.2), Inches(5.5), Inches(0.4), "老板 / 管理层", font_size=18, bold=True, color=ORANGE)
mgmt_items = ["数据分析：BI 导出 Excel → AI 分析异常", "决策推演：降价 10% → AI 推算利润影响", "行业研报速读：PDF → 提取品类趋势"]
for i, item in enumerate(mgmt_items):
    add_textbox(slide, Inches(7.2), Inches(5.65) + Inches(0.42)*i, Inches(5.5), Inches(0.4), f"▸  {item}", font_size=12, color=LIGHT)

add_bottom_bar(slide)
add_page_number(slide, 5)

# ═══════════════════════════════════
# SLIDE 6 - PROMOTION STRATEGY
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "推广方式：传染式扩散", "不做全员动员大会，让效果自己说话")

steps = [
    ("第 1-2 周", "选 2-3 个「种子选手」", "找跨境团队里最愿意尝试新事物的 1-2 人，私下教会他们用，让他们自己感受到效率提升。不搞全员培训。", ACCENT),
    ("第 3-4 周", "让效果自己说话", "种子选手在日常工作中用 AI 产出结果。安排一次非正式 10 分钟分享——让他们在部门周会上演示实际用法。", ACCENT2),
    ("第 2-3 月", "沉淀模板 + 按需扩散", "把已验证好用的 Prompt 整理成共享文档，按场景分类。谁感兴趣谁来拿，不强制推广。", GREEN),
    ("持续", "轻量激励", "每月评一个「AI 最佳实践」，奖金 500-1000 或老板在群里表扬。鼓励大家分享「今天用 AI 搞定了什么」。", ORANGE),
]

for i, (time, title, desc, color) in enumerate(steps):
    y = Inches(2.0) + Inches(1.25) * i

    # Timeline connector
    if i < len(steps) - 1:
        connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.45), y + Inches(0.65), Inches(0.04), Inches(0.95))
        connector.fill.solid()
        connector.fill.fore_color.rgb = RGBColor(0x2A, 0x3A, 0x6A)
        connector.line.fill.background()

    # Time badge
    badge = add_rect(slide, Inches(0.8), y + Inches(0.1), Inches(1.4), Inches(0.55), color, 0.12)
    add_textbox(slide, Inches(0.8), y + Inches(0.15), Inches(1.4), Inches(0.45), time, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.3), y + Inches(0.18), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_textbox(slide, Inches(2.3), y + Inches(0.2), Inches(0.35), Inches(0.3), str(i+1), font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Content
    add_textbox(slide, Inches(3.0), y + Inches(0.05), Inches(9), Inches(0.45), title, font_size=22, bold=True, color=WHITE)
    add_textbox(slide, Inches(3.0), y + Inches(0.5), Inches(9), Inches(0.55), desc, font_size=14, color=GRAY)

add_bottom_bar(slide)
add_page_number(slide, 6)

# ═══════════════════════════════════
# SLIDE 7 - RESOURCES
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "资源需求（极简版）", "全年总预算不到 3 万元")

# Budget table as cards
budget_items = [
    ("AI 工具账号", "ChatGPT Team / Claude Pro\n先给核心用户开 5-10 个账号", "1-2 万/年", ACCENT),
    ("Prompt 模板建设", "种子选手 + IT 一起整理\n沉淀到共享文档即可", "0 元", ACCENT2),
    ("内部推广", "奖励金 + 分享会茶歇\n月度 AI 最佳实践奖金", "0.5-1 万/年", GREEN),
]

for i, (title, desc, cost, color) in enumerate(budget_items):
    left = Inches(0.8) + Inches(4.1) * i
    card = add_rect(slide, left, Inches(2.2), Inches(3.7), Inches(2.8), BASE_CARD, 0.1)
    add_rect(slide, left, Inches(2.2), Inches(3.7), Inches(0.06), color)
    add_textbox(slide, left + Inches(0.3), Inches(2.5), Inches(3.1), Inches(0.4), title, font_size=20, bold=True, color=color)
    add_textbox(slide, left + Inches(0.3), Inches(3.0), Inches(3.1), Inches(1.0), desc, font_size=13, color=GRAY)
    add_textbox(slide, left + Inches(0.3), Inches(4.2), Inches(3.1), Inches(0.5), cost, font_size=28, bold=True, color=color)

# Total
add_rect(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(0.9), ACCENT, 0.12)
add_textbox(slide, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.7), "全年总预算 < 3 万元", font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.7), Inches(12), Inches(0.4),
            "对比：传统 AI 项目动辄 100-200 万起步。轻量方案用 1% 的成本，解决 80% 的日常效率问题。",
            font_size=13, color=SUBTLE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide)
add_page_number(slide, 7)

# ═══════════════════════════════════
# SLIDE 8 - PRECAUTIONS
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_section_title(slide, "注意事项", "五个关键提醒")

cautions = [
    ("数据安全底线", "敏感数据（客户名单、财务数据、成本价）脱敏后再输入 AI。「竞争对手想看的东西，就别直接贴进去」", ACCENT),
    ("不追求完美", "AI 输出是初稿不是终稿，人要审核。别因为偶尔出错就否定整个工具。80 分就够用了。", ACCENT2),
    ("不强制 KPI", "这不是考核指标，是给人配一个帮手。谁觉得有用谁用，强制推广反而招来反感。", GREEN),
    ("选好工具", "跨境场景用 Claude / ChatGPT（英文强）；国内场景补 Kimi（中文好、超长文档）。不要只用一个模型。", ORANGE),
    ("老板带头用", "老板自己开会用 AI 做纪要、用 AI 分析报表。比任何动员会都管用。", ACCENT3),
]

for i, (title, desc, color) in enumerate(cautions):
    y = Inches(2.0) + Inches(1.0) * i
    card = add_rect(slide, Inches(0.8), y, Inches(12), Inches(0.85), BASE_CARD, 0.08)
    # Left color strip
    add_rect(slide, Inches(0.8), y, Inches(0.08), Inches(0.85), color)

    add_textbox(slide, Inches(1.3), y + Inches(0.08), Inches(3), Inches(0.35), f"⚠ {title}", font_size=17, bold=True, color=color)
    add_textbox(slide, Inches(1.3), y + Inches(0.45), Inches(10.5), Inches(0.35), desc, font_size=13, color=GRAY)

add_bottom_bar(slide)
add_page_number(slide, 8)

# ═══════════════════════════════════
# SLIDE 9 - SUMMARY
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Big quote style
add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.5), "一张图总结", font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

summary_lines = [
    ("不是 AI 项目，是 AI 习惯。", True, 36, ACCENT),
    ("", False, 12, WHITE),
    ("不是技术变革，是工作方式升级。", True, 28, WHITE),
    ("", False, 12, WHITE),
    ("不建系统、不招团队、不花大钱。", True, 28, WHITE),
    ("", False, 12, WHITE),
    ("就是给每个人配一个好用的「实习生」，", True, 28, WHITE),
    ("让效果自己传染。", True, 36, ACCENT2),
]

y_pos = Inches(2.0)
for txt, bold, size, color in summary_lines:
    if txt:
        add_textbox(slide, Inches(1.5), y_pos, Inches(10), Inches(0.8), txt, font_size=size, bold=bold, color=color, alignment=PP_ALIGN.CENTER)
        y_pos += Inches(0.7)
    else:
        y_pos += Inches(0.1)

add_bottom_bar(slide)
add_page_number(slide, 9)

# ═══════════════════════════════════
# SLIDE 10 - THANK YOU
# ═══════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)

# Decorative circles
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4.5), Inches(6), Inches(6))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x13, 0x1F, 0x42)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-1.5), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x13, 0x1F, 0x42)
circle2.line.fill.background()

add_textbox(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(1.0), "谢谢", font_size=56, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0), Inches(3.6), Inches(13.333), Inches(0.6), "轻量化 AI 赋能，从第一个人开始", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.05), ACCENT)

add_textbox(slide, Inches(0), Inches(5.0), Inches(13.333), Inches(0.5), "有问题随时沟通", font_size=16, color=SUBTLE, alignment=PP_ALIGN.CENTER)

add_bottom_bar(slide, "轻量化 AI 赋能策略 | 卫生护理企业 | 2025年6月")

# ── Save ──
output_path = "/Users/admin/claude code project/AI赋能落地策略-轻量版.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
